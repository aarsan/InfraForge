"""Generate InfraForge 2-slide presentation deck for the GitHub Copilot SDK Challenge."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # Dark background
BG_CARD      = RGBColor(0x16, 0x1B, 0x22)   # Card background
ACCENT_BLUE  = RGBColor(0x58, 0xA6, 0xFF)   # GitHub blue accent
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)   # Green for approved/success
ACCENT_ORANGE= RGBColor(0xD2, 0x9E, 0x22)   # Orange accent
TEXT_PRIMARY  = RGBColor(0xF0, 0xF6, 0xFC)   # White text
TEXT_SECONDARY= RGBColor(0x8B, 0x94, 0x9E)   # Muted text
TEXT_MUTED    = RGBColor(0x6E, 0x76, 0x81)   # Even more muted
BORDER_COLOR = RGBColor(0x30, 0x36, 0x3D)   # Subtle border
PURPLE       = RGBColor(0xBC, 0x8C, 0xFF)   # Purple accent
RED_ACCENT   = RGBColor(0xFF, 0x7B, 0x72)   # Red for problems

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=14, color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=14, color=TEXT_PRIMARY, bold=False, space_before=Pt(4), space_after=Pt(2), font_name="Segoe UI", alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_icon_text(tf, icon, text, size=13, color=TEXT_PRIMARY, bold=False, space_before=Pt(6)):
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = f"{icon}  {text}"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Segoe UI"
    p.space_before = space_before
    p.space_after = Pt(2)
    return p


def add_bullet(tf, text, size=12, color=TEXT_PRIMARY, indent_level=0, space_before=Pt(3)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.level = indent_level
    p.space_before = space_before
    p.space_after = Pt(1)
    return p


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Business Value Proposition
# ════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
add_bg(slide1)

# ── Top accent bar ──
accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(4))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = ACCENT_BLUE
accent_bar.line.fill.background()

# ── Title area ──
tb = add_text_box(slide1, Inches(0.6), Inches(0.25), Inches(8), Inches(1.0))
set_text(tb.text_frame, "InfraForge", size=40, color=TEXT_PRIMARY, bold=True)
add_paragraph(tb.text_frame, "Agentic Self-Service Infrastructure Platform", size=18, color=ACCENT_BLUE, bold=False, space_before=Pt(2))

# ── Tagline ──
tb2 = add_text_box(slide1, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5))
set_text(tb2.text_frame, "Agents write the policies, generate the infrastructure code, test it, and deploy it \u2014 powered by the GitHub Copilot SDK",
         size=14, color=TEXT_SECONDARY)

# ── "The Problem" card ──
card_left = Inches(0.5)
card_top = Inches(1.9)
card_w = Inches(5.8)
card_h = Inches(2.4)
add_shape(slide1, card_left, card_top, card_w, card_h, BG_CARD, BORDER_COLOR, 0.02)

tb_prob = add_text_box(slide1, card_left + Inches(0.3), card_top + Inches(0.15), card_w - Inches(0.5), card_h - Inches(0.3))
set_text(tb_prob.text_frame, "The Problem", size=18, color=RED_ACCENT, bold=True)
add_icon_text(tb_prob.text_frame, "⏳", "App teams wait days-to-weeks for platform teams to write IaC", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_prob.text_frame, "🔄", "Every team reinvents patterns — inconsistent, insecure defaults", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_prob.text_frame, "💸", "Cost surprises after deployment, not before", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_prob.text_frame, "🧱", "No template reuse — written once, lost in repo sprawl", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_prob.text_frame, "🚧", "Governance bottleneck: security, compliance, policy reviews are manual", size=12, color=TEXT_PRIMARY)

# ── "The Solution" card ──
card2_left = Inches(6.6)
card2_w = Inches(6.2)
add_shape(slide1, card2_left, card_top, card2_w, card_h, BG_CARD, ACCENT_BLUE, 0.02)

tb_sol = add_text_box(slide1, card2_left + Inches(0.3), card_top + Inches(0.15), card2_w - Inches(0.5), card_h - Inches(0.3))
set_text(tb_sol.text_frame, "The InfraForge Solution", size=18, color=ACCENT_GREEN, bold=True)
add_icon_text(tb_sol.text_frame, "🤖", "Agents write IaC: ARM, Bicep & Terraform from natural language", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_sol.text_frame, "🛡️", "Agents author policies: security, governance & compliance auto-generated", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_sol.text_frame, "🧪", "Agents test & validate: deploy to Azure, run tests, self-heal on failure", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_sol.text_frame, "🚀", "Agents deploy: ARM SDK deployment with live progress streaming", size=12, color=TEXT_PRIMARY)
add_icon_text(tb_sol.text_frame, "📚", "Catalog-first: reuse approved templates, register new ones for the org", size=12, color=TEXT_PRIMARY)

# ── Key Value Metrics row ──
metrics = [
    ("Days → Minutes", "Infrastructure\nProvisioning", ACCENT_BLUE),
    ("100% Automated", "Governance &\nCompliance", ACCENT_GREEN),
    ("Zero Terraform\nExpertise Needed", "Self-Service\nfor App Teams", PURPLE),
    ("Full Cost\nTransparency", "Before\nDeployment", ACCENT_ORANGE),
]

metric_top = Inches(4.55)
metric_h = Inches(1.3)
metric_gap = Inches(0.25)
total_w = Inches(12.333)
metric_w = (total_w - metric_gap * (len(metrics) - 1)) / len(metrics)

for i, (title, subtitle, color) in enumerate(metrics):
    mx = Inches(0.5) + i * (metric_w + metric_gap)
    add_shape(slide1, mx, metric_top, metric_w, metric_h, BG_CARD, color, 0.03)

    mtb = add_text_box(slide1, mx + Inches(0.15), metric_top + Inches(0.12), metric_w - Inches(0.3), metric_h - Inches(0.2))
    set_text(mtb.text_frame, title, size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(mtb.text_frame, subtitle, size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER, space_before=Pt(6))

# ── Integrations / Tech bar ──
tech_top = Inches(6.1)
tech_h = Inches(0.7)
add_shape(slide1, Inches(0.5), tech_top, Inches(12.333), tech_h, BG_CARD, BORDER_COLOR, 0.02)

tb_tech = add_text_box(slide1, Inches(0.8), tech_top + Inches(0.1), Inches(12), tech_h - Inches(0.2))
set_text(tb_tech.text_frame, "Built With:  GitHub Copilot SDK  •  Microsoft Entra ID  •  Azure SQL  •  Azure ARM SDK  •  FastAPI  •  Microsoft Work IQ  •  Microsoft Fabric  •  Claude (Anthropic)",
         size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# ── Footer ──
tb_footer = add_text_box(slide1, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4))
set_text(tb_footer.text_frame, "GitHub Copilot SDK Enterprise Challenge  |  March 2026  |  github.com/aharsan/InfraForge",
         size=11, color=TEXT_MUTED, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture & Judging Criteria Alignment
# ════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

# ── Top accent bar ──
accent_bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(4))
accent_bar2.fill.solid()
accent_bar2.fill.fore_color.rgb = ACCENT_BLUE
accent_bar2.line.fill.background()

# ── Title ──
tb_title2 = add_text_box(slide2, Inches(0.6), Inches(0.25), Inches(10), Inches(0.7))
set_text(tb_title2.text_frame, "Architecture", size=34, color=TEXT_PRIMARY, bold=True)

# ──────────────────────────────────────────────────────────────────────
# Simplified architecture diagram (left side)
# ──────────────────────────────────────────────────────────────────────
diag_left  = Inches(0.4)
diag_top   = Inches(1.1)
diag_w     = Inches(6.0)
diag_h     = Inches(5.5)

add_shape(slide2, diag_left, diag_top, diag_w, diag_h, BG_CARD, BORDER_COLOR, 0.01)

def arch_box(slide, x, y, w, h, label, color, text_color=TEXT_PRIMARY, font_size=10):
    s = add_shape(slide, x, y, w, h, color, None, 0.04)
    s.text_frame.word_wrap = True
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = s.text_frame.paragraphs[0].add_run()
    r.text = label
    r.font.size = Pt(font_size)
    r.font.color.rgb = text_color
    r.font.bold = True
    r.font.name = "Segoe UI"
    return s

def arrow_down(slide, cx, y_start, y_end, color=ACCENT_BLUE):
    line = slide.shapes.add_connector(1, cx, y_start, cx, y_end)
    line.line.color.rgb = color
    line.line.width = Pt(2)

def arrow_right(slide, x_start, x_end, cy, color=ACCENT_BLUE):
    line = slide.shapes.add_connector(1, x_start, cy, x_end, cy)
    line.line.color.rgb = color
    line.line.width = Pt(2)

box_h = Inches(0.5)
pad = Inches(0.3)
inner_left = diag_left + pad
inner_w = diag_w - pad * 2

# Row 1 — Users
r1_y = diag_top + Inches(0.3)
lbl = add_text_box(slide2, inner_left, r1_y - Inches(0.15), Inches(1), Inches(0.2))
set_text(lbl.text_frame, "USERS", size=8, color=TEXT_MUTED, bold=True)

box_w3 = (inner_w - Inches(0.3)) / 3
arch_box(slide2, inner_left, r1_y + Inches(0.05), box_w3, box_h, "Web UI (SPA)", RGBColor(0x1F, 0x2A, 0x37), ACCENT_BLUE, 10)
arch_box(slide2, inner_left + box_w3 + Inches(0.15), r1_y + Inches(0.05), box_w3, box_h, "CLI", RGBColor(0x1F, 0x2A, 0x37), ACCENT_BLUE, 10)
arch_box(slide2, inner_left + (box_w3 + Inches(0.15)) * 2, r1_y + Inches(0.05), box_w3, box_h, "Entra ID (SSO)", RGBColor(0x1A, 0x25, 0x33), PURPLE, 10)

mid_x = inner_left + inner_w / 2
arrow_down(slide2, mid_x, r1_y + Inches(0.05) + box_h, r1_y + Inches(0.05) + box_h + Inches(0.2))

# Row 2 — Backend
r2_y = r1_y + box_h + Inches(0.3)
lbl2 = add_text_box(slide2, inner_left, r2_y - Inches(0.15), Inches(1.2), Inches(0.2))
set_text(lbl2.text_frame, "BACKEND", size=8, color=TEXT_MUTED, bold=True)
arch_box(slide2, inner_left, r2_y + Inches(0.05), inner_w, box_h,
         "FastAPI  +  WebSocket Chat  +  Pipeline Engine", RGBColor(0x15, 0x23, 0x30), ACCENT_GREEN, 11)
arrow_down(slide2, mid_x, r2_y + Inches(0.05) + box_h, r2_y + Inches(0.05) + box_h + Inches(0.2))

# Row 3 — AI Engine
r3_y = r2_y + box_h + Inches(0.3)
lbl3 = add_text_box(slide2, inner_left, r3_y - Inches(0.15), Inches(1.5), Inches(0.2))
set_text(lbl3.text_frame, "AGENTIC ENGINE", size=8, color=TEXT_MUTED, bold=True)
half_w = (inner_w - Inches(0.15)) / 2
arch_box(slide2, inner_left, r3_y + Inches(0.05), half_w, Inches(0.6),
         "GitHub Copilot SDK\n(Model Router)", RGBColor(0x1A, 0x15, 0x30), PURPLE, 10)
arch_box(slide2, inner_left + half_w + Inches(0.15), r3_y + Inches(0.05), half_w, Inches(0.6),
         "16 Agent Tools\n(Generate, Deploy, Govern, Cost...)", RGBColor(0x1A, 0x25, 0x15), ACCENT_GREEN, 10)
arrow_right(slide2, inner_left + half_w, inner_left + half_w + Inches(0.15),
            r3_y + Inches(0.05) + Inches(0.3), PURPLE)
arrow_down(slide2, mid_x, r3_y + Inches(0.65), r3_y + Inches(0.85))

# Row 4 — Data & Deploy
r4_y = r3_y + Inches(0.9)
lbl4 = add_text_box(slide2, inner_left, r4_y - Inches(0.15), Inches(1.5), Inches(0.2))
set_text(lbl4.text_frame, "DATA & DEPLOY", size=8, color=TEXT_MUTED, bold=True)
third_w = (inner_w - Inches(0.3)) / 3
arch_box(slide2, inner_left, r4_y + Inches(0.05), third_w, Inches(0.6),
         "Azure SQL\n(Catalog + Governance)", RGBColor(0x1F, 0x2A, 0x37), ACCENT_BLUE, 9)
arch_box(slide2, inner_left + third_w + Inches(0.15), r4_y + Inches(0.05), third_w, Inches(0.6),
         "ARM SDK\n(What-If + Deploy)", RGBColor(0x2A, 0x1F, 0x15), ACCENT_ORANGE, 9)
arch_box(slide2, inner_left + (third_w + Inches(0.15)) * 2, r4_y + Inches(0.05), third_w, Inches(0.6),
         "Microsoft Fabric\n(Analytics)", RGBColor(0x1A, 0x25, 0x15), ACCENT_GREEN, 9)

# Row 5 — Pipeline (simplified)
r5_y = r4_y + Inches(0.9)
lbl5 = add_text_box(slide2, inner_left, r5_y - Inches(0.15), Inches(2), Inches(0.2))
set_text(lbl5.text_frame, "AGENTIC PIPELINE", size=8, color=TEXT_MUTED, bold=True)
pipe_steps = ["Govern", "Generate\nARM", "Generate\nPolicy", "CISO\nReview", "Deploy\n& Heal", "Test &\nPromote"]
pipe_colors = [ACCENT_BLUE, PURPLE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_BLUE, ACCENT_GREEN]
step_w = (inner_w - Inches(0.1) * 5) / 6
for i, (step, clr) in enumerate(zip(pipe_steps, pipe_colors)):
    sx = inner_left + i * (step_w + Inches(0.1))
    arch_box(slide2, sx, r5_y + Inches(0.05), step_w, Inches(0.5), step, RGBColor(0x1F, 0x2A, 0x37), clr, 8)
    if i < len(pipe_steps) - 1:
        arrow_right(slide2, sx + step_w, sx + step_w + Inches(0.1),
                    r5_y + Inches(0.05) + Inches(0.25), clr)

# ──────────────────────────────────────────────────────────────────────
# Right side — Judging criteria cards (what the hackathon scores on)
# ──────────────────────────────────────────────────────────────────────
panel_left = Inches(6.7)
panel_w = Inches(6.2)
card_gap = Inches(0.2)

criteria = [
    ("Enterprise Value & Reusability", "30 pts", ACCENT_BLUE, [
        "Self-service for any team — zero IaC expertise needed",
        "Catalog-first: approved templates reused across the org",
        "Natural language intake replaces tickets & wait times",
    ]),
    ("Azure & Microsoft Integration", "25 pts", PURPLE, [
        "GitHub Copilot SDK — agentic AI engine with 16 tools",
        "Entra ID SSO — identity-aware provisioning & tagging",
        "Azure SQL, ARM SDK, Azure Policy, Microsoft Fabric",
    ]),
    ("Operational Readiness", "15 pts", ACCENT_GREEN, [
        "Self-healing deployment — auto-retry on ARM failures",
        "Live deploy progress streaming via WebSocket",
        "Full pipeline: govern → generate → test → promote",
    ]),
    ("Security, Governance & RAI", "15 pts", ACCENT_ORANGE, [
        "Governance-first: agents check approval before generating",
        "Virtual CISO reviews every onboarding for risk & compliance",
        "Policy engine enforces tags, naming, regions, SKUs",
    ]),
]

card_h = (Inches(5.5) - card_gap * (len(criteria) - 1)) / len(criteria)

for i, (title, pts, color, bullets) in enumerate(criteria):
    cy = diag_top + i * (card_h + card_gap)
    add_shape(slide2, panel_left, cy, panel_w, card_h, BG_CARD, color, 0.02)

    # Title row with points badge
    tb_title = add_text_box(slide2, panel_left + Inches(0.25), cy + Inches(0.1), panel_w - Inches(1.3), Inches(0.35))
    set_text(tb_title.text_frame, title, size=15, color=color, bold=True)

    tb_pts = add_text_box(slide2, panel_left + panel_w - Inches(1.0), cy + Inches(0.1), Inches(0.8), Inches(0.3))
    set_text(tb_pts.text_frame, pts, size=11, color=TEXT_MUTED, bold=True, alignment=PP_ALIGN.RIGHT)

    # Bullet points
    tb_bullets = add_text_box(slide2, panel_left + Inches(0.25), cy + Inches(0.4), panel_w - Inches(0.5), card_h - Inches(0.5))
    first = True
    for bullet in bullets:
        if first:
            set_text(tb_bullets.text_frame, f"•  {bullet}", size=11, color=TEXT_PRIMARY)
            first = False
        else:
            add_paragraph(tb_bullets.text_frame, f"•  {bullet}", size=11, color=TEXT_PRIMARY, space_before=Pt(3))

# ── Bonus bar ──
bonus_top = diag_top + Inches(5.7)
bonus_w = Inches(12.333)
add_shape(slide2, Inches(0.4), bonus_top, bonus_w, Inches(0.55), BG_CARD, BORDER_COLOR, 0.02)
tb_bonus = add_text_box(slide2, Inches(0.7), bonus_top + Inches(0.08), bonus_w - Inches(0.6), Inches(0.4))
set_text(tb_bonus.text_frame, "Bonus:   Microsoft Work IQ (+15 pts) — M365 org intelligence via MCP (emails, meetings, docs, Teams, people)   •   Fabric IQ (+15 pts) — usage analytics synced to Fabric",
         size=12, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# ── Footer ──
tb_footer2 = add_text_box(slide2, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4))
set_text(tb_footer2.text_frame, "Repo: github.com/aharsan/InfraForge  |  GitHub Copilot SDK Enterprise Challenge  |  March 2026",
         size=11, color=TEXT_MUTED, alignment=PP_ALIGN.LEFT)


# ── Save ──
out_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "presentations", "InfraForge.pptx")
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"Saved: {out_path}")
